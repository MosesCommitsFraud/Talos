import asyncio
import base64
import hashlib
import json
import mimetypes
import os
import shlex
import shutil
import signal
import subprocess
import threading
import time
import uuid
import fnmatch
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from pydantic import BaseModel


HOME_ROOT = Path(os.getenv("TALOS_SANDBOX_HOME_ROOT", "/home/talos"))
STATE_PATH = Path(os.getenv("TALOS_SANDBOX_STATE", "/var/lib/talos-sandbox/state.json"))

app = FastAPI(title="Talos Sandbox", version="0.1.0")

# Shared-secret auth: when TALOS_SANDBOX_KEY is set, every request (except the
# health probe) must carry it in the X-Talos-Sandbox-Key header. The app sends it
# on every call. Unset = auth disabled (dev/local only) — a warning is logged.
SANDBOX_KEY = os.getenv("TALOS_SANDBOX_KEY", "").strip()


@app.middleware("http")
async def _require_sandbox_key(request: Request, call_next):
    if SANDBOX_KEY and request.url.path != "/health":
        if request.headers.get("x-talos-sandbox-key", "") != SANDBOX_KEY:
            from fastapi.responses import JSONResponse
            return JSONResponse({"detail": "unauthorized"}, status_code=401)
    return await call_next(request)


if not SANDBOX_KEY:
    import logging as _logging
    _logging.getLogger("uvicorn.error").warning(
        "TALOS_SANDBOX_KEY is not set — sandbox API is UNAUTHENTICATED. Set it in .env to require auth."
    )


class EnsureUserResponse(BaseModel):
    user_id: str
    linux_user: str
    home: str


class WorkspaceResponse(BaseModel):
    user_id: str
    chat_id: str
    linux_user: str
    workspace: str
    filename: str | None = None
    path: str | None = None


class ExecRequest(BaseModel):
    kind: str = "bash"
    command: str = ""
    code: str = ""
    timeout: int = 120


class ExecResponse(BaseModel):
    user_id: str
    chat_id: str
    linux_user: str
    workspace: str
    stdout: str
    stderr: str
    exit_code: int
    timed_out: bool = False
    # Image files created/modified by the run (matplotlib plots, etc.), returned
    # as base64 data URLs so the chat can display them inline.
    images: list[dict[str, str]] = []
    image_note: str = ""
    created_artifacts: list[str] = []


class ExecuteRequest(BaseModel):
    command: str
    cwd: str | None = None
    env: dict[str, str] | None = None


class InputRequest(BaseModel):
    input: str


class CwdRequest(BaseModel):
    path: str


class FileReadRequest(BaseModel):
    path: str
    offset: int = 0
    limit: int = 0


class FileWriteRequest(BaseModel):
    path: str
    content: str = ""


class FileEditRequest(BaseModel):
    path: str
    old_string: str = ""
    new_string: str = ""
    replace_all: bool = False
    # OpenTerminal-style multi-chunk edit. Each item: {target, replacement,
    # start_line?, end_line?, allow_multiple?}. Applied in order; takes
    # precedence over old_string/new_string when non-empty.
    edits: list[dict[str, Any]] = []


class PathRequest(BaseModel):
    path: str


class MoveRequest(BaseModel):
    src: str
    dst: str


class CellRequest(BaseModel):
    code: str = ""
    timeout: int = 0


class SearchRequest(BaseModel):
    pattern: str = ""
    path: str = ""
    glob: str = ""
    ignore_case: bool = False
    max_results: int = 200


class ListRequest(BaseModel):
    path: str = ""


MAX_READ_CHARS = 20_000
MAX_OUTPUT_CHARS = 10_000
MAX_DIFF_LINES = 400
SKIP_DIRS = {".git", ".talos-home", ".talos-preview", "node_modules", "__pycache__", ".venv", "venv", "dist", "build", ".pytest_cache", ".mypy_cache"}
PROCESS_LOG_ROOT = STATE_PATH.parent / "processes"
PROCESS_RETENTION_SECONDS = 3600
SESSION_CWD_TTL_SECONDS = int(os.getenv("TALOS_SANDBOX_SESSION_CWD_TTL", "604800"))

# Images a run wants to PRESENT to the user (final charts/results) are returned
# inline as base64 data URLs. To opt in, the code must save the file under an
# `output/` directory in the workspace — scratch/WIP images written anywhere else
# are ignored. Caps keep session history from bloating; typical plots are tens of
# KB, so these limits only trip on pathological output.
OUTPUT_DIR_NAME = "output"
IMAGE_MIME_BY_EXT = {
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp",
    ".svg": "image/svg+xml",
}
MAX_IMAGE_BYTES = 3_000_000        # per file
MAX_IMAGES = 6                     # per run
MAX_IMAGES_TOTAL_BYTES = 9_000_000  # combined

# Files that are build/runtime noise, not deliverables — hidden from the artifacts list.
ARTIFACT_JUNK_EXTS = {".pyc", ".pyo", ".pyd", ".class", ".o", ".obj"}
ARTIFACT_JUNK_NAMES = {".DS_Store", "Thumbs.db", ".gitignore", ".python-version"}

# Binary document formats that read_file extracts to text instead of returning raw bytes.
DOC_EXTRACT_EXTS = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}


@dataclass
class BackgroundProcess:
    id: str
    user_id: str
    chat_id: str
    command: str
    runner: "PtyRunner"
    log_path: Path
    status: str = "running"
    exit_code: int | None = None
    started_at: float = field(default_factory=time.time)
    finished_at: float | None = None
    log_task: asyncio.Task | None = field(default=None, repr=False)


_processes: dict[str, BackgroundProcess] = {}
_session_cwds: dict[str, tuple[str, float]] = {}
_repaired_workspaces: set[str] = set()
_workspace_repair_lock = threading.Lock()
# Persistent per-workspace Python kernels: session_key -> {"proc", "sock"}.
_kernels: dict[str, dict[str, Any]] = {}

# Embedded "kernel server": runs AS THE SANDBOX USER (via gosu) and keeps a
# persistent namespace between calls. Talks length-prefixed JSON over a Unix
# socket in the workspace. argv: [sock_path, workdir].
_KERNEL_SERVER_SRC = r'''
import socket, sys, json, struct, io, contextlib, traceback, os, ast
sock_path, workdir = sys.argv[1], sys.argv[2]
try:
    os.chdir(workdir)
except Exception:
    pass
ns = {"__name__": "__main__"}
try:
    if os.path.exists(sock_path):
        os.unlink(sock_path)
except OSError:
    pass
srv = socket.socket(socket.AF_UNIX, socket.SOCK_STREAM)
srv.bind(sock_path)
try:
    os.chmod(sock_path, 0o666)
except OSError:
    pass
srv.listen(1)
def _recvall(c, n):
    buf = b""
    while len(buf) < n:
        d = c.recv(n - len(buf))
        if not d:
            return None
        buf += d
    return buf
def _read(c):
    h = _recvall(c, 4)
    if not h:
        return None
    (ln,) = struct.unpack(">I", h)
    b = _recvall(c, ln)
    return json.loads(b.decode("utf-8")) if b is not None else None
def _send(c, obj):
    data = json.dumps(obj).encode("utf-8")
    c.sendall(struct.pack(">I", len(data)) + data)
while True:
    try:
        conn, _ = srv.accept()
    except Exception:
        break
    try:
        msg = _read(conn)
        if msg is None:
            conn.close()
            continue
        code = msg.get("code", "")
        out, err, error = io.StringIO(), io.StringIO(), ""
        try:
            with contextlib.redirect_stdout(out), contextlib.redirect_stderr(err):
                try:
                    tree = ast.parse(code, "<cell>", "exec")
                except SyntaxError:
                    exec(compile(code, "<cell>", "exec"), ns)
                else:
                    if tree.body and isinstance(tree.body[-1], ast.Expr):
                        last = ast.Expression(tree.body.pop().value)
                        ast.fix_missing_locations(last)
                        if tree.body:
                            exec(compile(tree, "<cell>", "exec"), ns)
                        val = eval(compile(last, "<cell>", "eval"), ns)
                        if val is not None:
                            print(repr(val))
                    else:
                        exec(compile(tree, "<cell>", "exec"), ns)
        except Exception:
            error = traceback.format_exc()
        _send(conn, {"stdout": out.getvalue(), "stderr": err.getvalue(), "error": error})
    except Exception:
        pass
    finally:
        try:
            conn.close()
        except Exception:
            pass
'''


def linux_user(user_id: str) -> str:
    digest = hashlib.sha256(user_id.encode("utf-8")).hexdigest()[:16]
    return f"talos_{digest}"


def user_home(user_id: str) -> Path:
    return HOME_ROOT / linux_user(user_id)


def workspace_path(user_id: str, chat_id: str) -> Path:
    safe_chat = hashlib.sha256(chat_id.encode("utf-8")).hexdigest()[:24]
    return user_home(user_id) / "workspaces" / safe_chat


def _session_key(user_id: str, chat_id: str) -> str:
    return f"{user_id}:{chat_id}"


def _workspace(user_id: str, chat_id: str) -> tuple[str, Path]:
    name, home = ensure_user(user_id)
    if home.is_symlink():
        raise HTTPException(500, "Unsafe sandbox home")
    workspaces = home / "workspaces"
    if workspaces.is_symlink():
        raise HTTPException(500, "Unsafe sandbox workspace root")
    workspaces.mkdir(parents=True, exist_ok=True)
    workspace = workspace_path(user_id, chat_id)
    if workspace.is_symlink():
        raise HTTPException(500, "Unsafe sandbox workspace")
    workspace.mkdir(parents=True, exist_ok=True)
    # The path above each workspace stays root-managed so the sandbox account
    # cannot replace the workspace root with a symlink before root repairs it.
    _run(["chown", "root:root", str(home), str(workspaces)])
    home.chmod(0o711)
    workspaces.chmod(0o711)
    _run(["chown", f"{name}:{name}", str(workspace)])
    private_home = workspace / ".talos-home"
    private_home.mkdir(exist_ok=True)
    _run(["chown", f"{name}:{name}", str(private_home)])
    private_home.chmod(0o700)

    # Migrate legacy/restored files created by the root-running daemon. Tool
    # processes run as the per-user account and must be able to update every
    # listed artifact and create sibling temp files for atomic Office saves.
    repair_key = str(workspace.resolve())
    with _workspace_repair_lock:
        if repair_key not in _repaired_workspaces:
            _repair_workspace_tree(name, workspace)
            _repaired_workspaces.add(repair_key)
    return name, workspace


def _repair_workspace_tree(name: str, root: Path) -> None:
    """Descriptor-relative repair that cannot be redirected through symlinks."""
    import pwd
    import stat

    account = pwd.getpwnam(name)
    if root.is_symlink() or not root.is_dir():
        raise HTTPException(500, "Unsafe sandbox workspace")
    for _current, dirs, files, dir_fd in os.fwalk(root, topdown=True, follow_symlinks=False):
        os.fchown(dir_fd, account.pw_uid, account.pw_gid)
        os.fchmod(dir_fd, os.fstat(dir_fd).st_mode | 0o700)
        all_dirs = list(dirs)
        dirs[:] = [entry for entry in dirs if entry not in SKIP_DIRS]
        for entry_name in all_dirs + files:
            try:
                entry_stat = os.stat(entry_name, dir_fd=dir_fd, follow_symlinks=False)
                os.chown(
                    entry_name,
                    account.pw_uid,
                    account.pw_gid,
                    dir_fd=dir_fd,
                    follow_symlinks=False,
                )
                if stat.S_ISLNK(entry_stat.st_mode):
                    continue
                flags = os.O_RDONLY | os.O_NONBLOCK | os.O_NOFOLLOW | os.O_CLOEXEC
                if stat.S_ISDIR(entry_stat.st_mode):
                    flags |= os.O_DIRECTORY
                entry_fd = os.open(entry_name, flags, dir_fd=dir_fd)
                try:
                    opened_stat = os.fstat(entry_fd)
                    os.fchown(entry_fd, account.pw_uid, account.pw_gid)
                    os.fchmod(
                        entry_fd,
                        opened_stat.st_mode | (0o700 if stat.S_ISDIR(opened_stat.st_mode) else 0o600),
                    )
                finally:
                    os.close(entry_fd)
            except OSError:
                continue


def _chown_user_chain(name: str, workspace: Path, path: Path) -> None:
    """chown `path` and each ancestor directory up to the workspace root to the
    sandbox user. The daemon runs as root, so dirs it creates via mkdir would be
    root-owned and lock the user (gosu) out of writing into them later."""
    try:
        root = workspace.resolve()
        node = path.resolve()
    except OSError:
        return
    if node != root and root not in node.parents:
        return
    targets: list[Path] = []
    while True:
        targets.append(node)
        if node == root or root not in node.parents:
            break
        node = node.parent
    for t in targets:
        subprocess.run(["chown", f"{name}:{name}", str(t)], capture_output=True)


def _safe_path(workspace: Path, raw_path: str) -> Path:
    raw = (raw_path or ".").strip() or "."
    candidate = Path(raw)
    if candidate.is_absolute():
        candidate = candidate
    else:
        candidate = workspace / candidate
    resolved = candidate.resolve()
    root = workspace.resolve()
    if resolved != root and root not in resolved.parents:
        raise HTTPException(403, "Path escapes workspace")
    return resolved


def _expire_session_cwds() -> None:
    now = time.time()
    for key, (_cwd, ts) in list(_session_cwds.items()):
        if now - ts > SESSION_CWD_TTL_SECONDS:
            _session_cwds.pop(key, None)


def _get_session_cwd(user_id: str, chat_id: str, workspace: Path) -> Path:
    _expire_session_cwds()
    key = _session_key(user_id, chat_id)
    cwd, _ts = _session_cwds.get(key, (str(workspace), 0))
    try:
        resolved = _safe_path(workspace, cwd)
    except HTTPException:
        resolved = workspace.resolve()
    _session_cwds[key] = (str(resolved), time.time())
    return resolved


def _set_session_cwd(user_id: str, chat_id: str, workspace: Path, raw_path: str) -> Path:
    target = _safe_path(workspace, raw_path)
    if not target.is_dir():
        raise HTTPException(404, "Directory not found")
    _session_cwds[_session_key(user_id, chat_id)] = (str(target), time.time())
    return target


def _truncate(text: str, limit: int = MAX_OUTPUT_CHARS) -> str:
    return text if len(text) <= limit else text[:limit] + f"\n... [truncated at {limit} chars]"


class PtyRunner:
    def __init__(self, *, linux_user: str, command: str, cwd: Path, env: dict[str, str] | None = None):
        import fcntl
        import pty
        import struct
        import termios

        master_fd, slave_fd = pty.openpty()
        try:
            fcntl.ioctl(slave_fd, termios.TIOCSWINSZ, struct.pack("HHHH", 24, 120, 0, 0))
            proc_env = os.environ.copy()
            if env:
                proc_env.update({str(k): str(v) for k, v in env.items()})
            if not env or "HOME" not in env:
                proc_env["HOME"] = str(HOME_ROOT / linux_user)
            proc_env["TERM"] = "xterm-256color"
            proc_env["COLUMNS"] = "120"
            proc_env["LINES"] = "40"
            proc_env["PATH"] = f"/opt/talos-sandbox-venv/bin:{proc_env.get('PATH', '')}"
            self.process = subprocess.Popen(
                ["gosu", linux_user, "bash", "-lc", command],
                stdin=slave_fd,
                stdout=slave_fd,
                stderr=slave_fd,
                cwd=str(cwd),
                env=proc_env,
                start_new_session=True,
            )
        except Exception:
            os.close(slave_fd)
            os.close(master_fd)
            raise
        os.close(slave_fd)
        self.master_fd = master_fd

    async def log_output(self, proc: BackgroundProcess) -> None:
        loop = asyncio.get_event_loop()
        proc.log_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            with proc.log_path.open("a", encoding="utf-8") as log:
                while True:
                    try:
                        data = await loop.run_in_executor(None, os.read, self.master_fd, 4096)
                    except OSError:
                        break
                    if not data:
                        break
                    log.write(json.dumps({"type": "output", "data": data.decode(errors="replace"), "ts": time.time()}) + "\n")
                    log.flush()
        finally:
            proc.exit_code = await asyncio.to_thread(self.process.wait)
            proc.status = "done"
            proc.finished_at = time.time()
            self.close()

    def write_input(self, text: str) -> None:
        os.write(self.master_fd, text.encode())

    def kill(self, force: bool = False) -> None:
        sig = signal.SIGKILL if force else signal.SIGTERM
        try:
            os.killpg(self.process.pid, sig)
        except (ProcessLookupError, PermissionError):
            pass

    def close(self) -> None:
        try:
            os.close(self.master_fd)
        except OSError:
            pass


def _read_process_log(path: Path, offset: int = 0, tail: int | None = None) -> tuple[str, int, bool]:
    if not path.exists():
        return "", 0, False
    lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    total = len(lines)
    selected = lines[max(offset, 0):]
    if tail and tail > 0:
        selected = selected[-tail:]
    chunks: list[str] = []
    truncated = False
    budget = MAX_OUTPUT_CHARS
    for line in selected:
        try:
            item = json.loads(line)
            text = str(item.get("data") or "")
        except Exception:
            text = line
        if len(text) > budget:
            chunks.append(text[:budget])
            truncated = True
            break
        chunks.append(text)
        budget -= len(text)
        if budget <= 0:
            truncated = True
            break
    return "".join(chunks), total, truncated


def _cleanup_processes() -> None:
    now = time.time()
    for pid, proc in list(_processes.items()):
        if proc.finished_at and now - proc.finished_at > PROCESS_RETENTION_SECONDS:
            _processes.pop(pid, None)


def _diff(old: str, new: str, path: str) -> dict[str, Any] | None:
    if old == new:
        return None
    import difflib
    lines = list(difflib.unified_diff(old.splitlines(), new.splitlines(), fromfile=f"a/{path}", tofile=f"b/{path}", lineterm=""))
    added = sum(1 for line in lines if line.startswith("+") and not line.startswith("+++"))
    removed = sum(1 for line in lines if line.startswith("-") and not line.startswith("---"))
    if len(lines) > MAX_DIFF_LINES:
        lines = lines[:MAX_DIFF_LINES] + [f"... diff truncated at {MAX_DIFF_LINES} lines"]
    return {"text": "\n".join(lines), "added": added, "removed": removed, "new_file": old == "", "file": Path(path).name or path}


def _run(args: list[str]) -> None:
    result = subprocess.run(args, text=True, capture_output=True)
    if result.returncode != 0:
        raise HTTPException(500, f"Command failed: {' '.join(args)}\n{result.stderr}")


def ensure_user(user_id: str) -> tuple[str, Path]:
    name = linux_user(user_id)
    home = user_home(user_id)
    HOME_ROOT.mkdir(parents=True, exist_ok=True)
    created = False
    if subprocess.run(["id", "-u", name], capture_output=True).returncode != 0:
        _run(["useradd", "--create-home", "--home-dir", str(home), "--shell", "/bin/bash", name])
        created = True
    if not home.exists():
        home.mkdir(parents=True, exist_ok=True)
        created = True
    # Workspace roots are daemon-managed; tool processes receive a writable
    # private HOME inside their own workspace instead.
    _run(["chown", "root:root", str(home)])
    home.chmod(0o711)
    return name, home


@app.get("/health")
def health() -> dict[str, Any]:
    return {"ok": True, "active": len(_processes)}


@app.post("/users/{user_id}/ensure", response_model=EnsureUserResponse)
def ensure_user_route(user_id: str) -> EnsureUserResponse:
    name, home = ensure_user(user_id)
    return EnsureUserResponse(user_id=user_id, linux_user=name, home=str(home))


@app.post("/users/{user_id}/workspaces/{chat_id}/ensure", response_model=WorkspaceResponse)
def ensure_workspace_route(user_id: str, chat_id: str) -> WorkspaceResponse:
    name, _home = ensure_user(user_id)
    workspace = workspace_path(user_id, chat_id)
    workspace.mkdir(parents=True, exist_ok=True)
    _run(["chown", "-R", f"{name}:{name}", str(workspace)])
    return WorkspaceResponse(user_id=user_id, chat_id=chat_id, linux_user=name, workspace=str(workspace))


@app.delete("/users/{user_id}/workspaces/{chat_id}")
def delete_workspace_route(user_id: str, chat_id: str) -> dict[str, Any]:
    """Remove a chat's workspace and everything in it. Called when the chat is
    deleted so files don't outlive the conversation. Idempotent — a missing
    workspace is a no-op."""
    # Kill any persistent kernel for this chat before removing its files.
    _stop_kernel(user_id, chat_id)
    workspace = workspace_path(user_id, chat_id)
    existed = workspace.exists()
    if existed:
        shutil.rmtree(workspace, ignore_errors=True)
    _repaired_workspaces.discard(str(workspace.resolve()))
    # Drop any cached session cwd for this chat so a recreated workspace starts clean.
    _session_cwds.pop(_session_key(user_id, chat_id), None)
    return {"ok": True, "deleted": existed, "workspace": str(workspace)}


@app.post("/users/{user_id}/workspaces/{chat_id}/upload", response_model=WorkspaceResponse)
async def upload_file_route(user_id: str, chat_id: str, file: UploadFile = File(...)) -> WorkspaceResponse:
    # Use _workspace so the workspace dir is owned by the sandbox user — otherwise
    # an upload would leave it root-owned and the agent couldn't write files later.
    name, workspace = _workspace(user_id, chat_id)
    filename = Path(file.filename or "upload.bin").name.replace("/", "_").replace("\\", "_")
    target = workspace / filename
    with target.open("wb") as out:
        shutil.copyfileobj(file.file, out)
    _run(["chown", f"{name}:{name}", str(target)])
    target.chmod(target.stat().st_mode | 0o600)
    return WorkspaceResponse(user_id=user_id, chat_id=chat_id, linux_user=name, workspace=str(workspace), filename=filename, path=str(target))


@app.get("/users/{user_id}/workspaces/{chat_id}/cwd")
def get_cwd_route(user_id: str, chat_id: str) -> dict[str, Any]:
    _name, workspace = _workspace(user_id, chat_id)
    cwd = _get_session_cwd(user_id, chat_id, workspace)
    return {"cwd": str(cwd), "workspace": str(workspace)}


@app.post("/users/{user_id}/workspaces/{chat_id}/cwd")
def set_cwd_route(user_id: str, chat_id: str, req: CwdRequest) -> dict[str, Any]:
    _name, workspace = _workspace(user_id, chat_id)
    cwd = _set_session_cwd(user_id, chat_id, workspace, req.path)
    return {"cwd": str(cwd), "workspace": str(workspace)}


async def _start_process(user_id: str, chat_id: str, req: ExecuteRequest) -> BackgroundProcess:
    name, workspace = _workspace(user_id, chat_id)
    base_cwd = _get_session_cwd(user_id, chat_id, workspace)
    cwd = _safe_path(workspace, req.cwd) if req.cwd else base_cwd
    if not cwd.is_dir():
        raise HTTPException(404, "Working directory not found")
    process_id = time.strftime("%Y%m%d-%H%M%S-") + uuid.uuid4().hex[:6]
    log_path = PROCESS_LOG_ROOT / f"{process_id}.jsonl"
    process_env = {**(req.env or {}), "HOME": str(workspace / ".talos-home")}
    runner = PtyRunner(linux_user=name, command=req.command, cwd=cwd, env=process_env)
    proc = BackgroundProcess(
        id=process_id,
        user_id=user_id,
        chat_id=chat_id,
        command=req.command,
        runner=runner,
        log_path=log_path,
    )
    proc.log_task = asyncio.create_task(runner.log_output(proc))
    _processes[process_id] = proc
    return proc


@app.get("/users/{user_id}/workspaces/{chat_id}/execute")
def list_processes_route(user_id: str, chat_id: str) -> list[dict[str, Any]]:
    _cleanup_processes()
    return [
        {
            "id": p.id,
            "command": p.command,
            "status": p.status,
            "exit_code": p.exit_code,
            "log_path": str(p.log_path),
        }
        for p in _processes.values()
        if p.user_id == user_id and p.chat_id == chat_id
    ]


@app.post("/users/{user_id}/workspaces/{chat_id}/execute")
async def execute_route(
    user_id: str,
    chat_id: str,
    req: ExecuteRequest,
    wait: float | None = None,
    tail: int | None = None,
) -> dict[str, Any]:
    proc = await _start_process(user_id, chat_id, req)
    if wait is not None:
        try:
            await asyncio.wait_for(asyncio.shield(proc.log_task), timeout=max(0, min(float(wait), 300)))
        except asyncio.TimeoutError:
            pass
    output, next_offset, truncated = _read_process_log(proc.log_path, offset=0, tail=tail)
    return {
        "id": proc.id,
        "command": proc.command,
        "status": proc.status,
        "exit_code": proc.exit_code,
        "output": output,
        "truncated": truncated,
        "next_offset": next_offset,
        "log_path": str(proc.log_path),
    }


def _get_process(user_id: str, chat_id: str, process_id: str) -> BackgroundProcess:
    _cleanup_processes()
    proc = _processes.get(process_id)
    if not proc or proc.user_id != user_id or proc.chat_id != chat_id:
        raise HTTPException(404, "Process not found")
    return proc


@app.get("/users/{user_id}/workspaces/{chat_id}/execute/{process_id}/status")
async def process_status_route(
    user_id: str,
    chat_id: str,
    process_id: str,
    wait: float | None = None,
    offset: int = 0,
    tail: int | None = None,
) -> dict[str, Any]:
    proc = _get_process(user_id, chat_id, process_id)
    if wait is not None and proc.status == "running":
        try:
            await asyncio.wait_for(asyncio.shield(proc.log_task), timeout=max(0, min(float(wait), 300)))
        except asyncio.TimeoutError:
            pass
    output, next_offset, truncated = _read_process_log(proc.log_path, offset=offset, tail=tail)
    return {
        "id": proc.id,
        "command": proc.command,
        "status": proc.status,
        "exit_code": proc.exit_code,
        "output": output,
        "truncated": truncated,
        "next_offset": next_offset,
        "log_path": str(proc.log_path),
    }


@app.post("/users/{user_id}/workspaces/{chat_id}/execute/{process_id}/input")
def process_input_route(user_id: str, chat_id: str, process_id: str, req: InputRequest) -> dict[str, Any]:
    proc = _get_process(user_id, chat_id, process_id)
    if proc.status != "running":
        raise HTTPException(400, "Process has already exited")
    proc.runner.write_input(req.input.encode("raw_unicode_escape").decode("unicode_escape"))
    return {"status": "ok"}


@app.delete("/users/{user_id}/workspaces/{chat_id}/execute/{process_id}")
def kill_process_route(user_id: str, chat_id: str, process_id: str, force: bool = False) -> dict[str, Any]:
    proc = _get_process(user_id, chat_id, process_id)
    if proc.status == "running":
        proc.runner.kill(force=force)
        proc.status = "killed"
        proc.finished_at = time.time()
    return {"status": proc.status}


@app.post("/users/{user_id}/workspaces/{chat_id}/terminals")
async def create_terminal_route(user_id: str, chat_id: str) -> dict[str, Any]:
    proc = await _start_process(user_id, chat_id, ExecuteRequest(command="bash"))
    return {"id": proc.id, "status": proc.status, "log_path": str(proc.log_path)}


@app.get("/users/{user_id}/workspaces/{chat_id}/terminals/{terminal_id}")
async def terminal_status_route(
    user_id: str,
    chat_id: str,
    terminal_id: str,
    offset: int = 0,
    tail: int | None = None,
) -> dict[str, Any]:
    return await process_status_route(user_id, chat_id, terminal_id, wait=None, offset=offset, tail=tail)


@app.post("/users/{user_id}/workspaces/{chat_id}/terminals/{terminal_id}/input")
def terminal_input_route(user_id: str, chat_id: str, terminal_id: str, req: InputRequest) -> dict[str, Any]:
    return process_input_route(user_id, chat_id, terminal_id, req)


@app.delete("/users/{user_id}/workspaces/{chat_id}/terminals/{terminal_id}")
def delete_terminal_route(user_id: str, chat_id: str, terminal_id: str, force: bool = False) -> dict[str, Any]:
    return kill_process_route(user_id, chat_id, terminal_id, force=force)


@app.get("/users/{user_id}/workspaces/{chat_id}/ports")
def list_ports_route(user_id: str, chat_id: str) -> dict[str, Any]:
    name, _workspace = _workspace(user_id, chat_id)
    try:
        result = subprocess.run(
            ["lsof", "-Pan", "-u", name, "-iTCP", "-sTCP:LISTEN"],
            text=True,
            capture_output=True,
            timeout=5,
        )
    except Exception as exc:
        return {"ports": [], "error": str(exc)}
    ports: list[dict[str, Any]] = []
    for line in (result.stdout or "").splitlines()[1:]:
        parts = line.split()
        if len(parts) < 9:
            continue
        name_part = parts[-2] if parts[-1] == "(LISTEN)" else parts[-1]
        match = re.search(r":(\d+)$", name_part)
        if not match:
            continue
        try:
            port = int(match.group(1))
        except ValueError:
            continue
        ports.append({"command": parts[0], "pid": parts[1], "port": port, "address": name_part})
    unique = {p["port"]: p for p in ports}
    return {"ports": list(unique.values())}


@app.api_route("/users/{user_id}/workspaces/{chat_id}/proxy/{port}/{path:path}", methods=["GET", "POST", "PUT", "PATCH", "DELETE", "HEAD", "OPTIONS"])
async def proxy_port_route(user_id: str, chat_id: str, port: int, path: str, request: Request):
    import httpx
    from fastapi.responses import Response

    if port < 1 or port > 65535:
        raise HTTPException(422, "Port must be between 1 and 65535")
    listed = list_ports_route(user_id, chat_id).get("ports", [])
    if port not in {int(p.get("port")) for p in listed}:
        raise HTTPException(404, "Port is not listening for this sandbox user")
    target = f"http://127.0.0.1:{port}/{path}"
    if request.query_params:
        target += f"?{request.query_params}"
    headers = dict(request.headers)
    for header in ("host", "authorization", "connection", "transfer-encoding"):
        headers.pop(header, None)
    async with httpx.AsyncClient(timeout=httpx.Timeout(120.0, connect=5.0), follow_redirects=False) as client:
        upstream = await client.request(request.method, target, headers=headers, content=await request.body() or None)
    response_headers = dict(upstream.headers)
    for header in ("transfer-encoding", "connection", "content-encoding", "content-length"):
        response_headers.pop(header, None)
    return Response(content=upstream.content, status_code=upstream.status_code, headers=response_headers)


def _collect_new_images(workspace: Path, since: float) -> tuple[list[dict[str, str]], str]:
    """Return images the run explicitly saved under an `output/` directory during
    this run (mtime >= since), as base64 data URLs, oldest-first so display order
    matches creation order. Images written anywhere else are treated as scratch /
    WIP and ignored. Returns (images, note) where note flags anything skipped."""
    candidates: list[tuple[float, int, Path, str]] = []
    try:
        for p in workspace.rglob("*"):
            ext = p.suffix.lower()
            if ext not in IMAGE_MIME_BY_EXT or not p.is_file():
                continue
            rel_parts = p.relative_to(workspace).parts
            if set(rel_parts) & SKIP_DIRS:
                continue
            # Opt-in: only files inside an `output/` directory are presented.
            if OUTPUT_DIR_NAME not in rel_parts[:-1]:
                continue
            try:
                st = p.stat()
            except OSError:
                continue
            # Small skew tolerance so a savefig flushed at run start still counts.
            if st.st_mtime + 0.5 < since:
                continue
            candidates.append((st.st_mtime, st.st_size, p, ext))
    except OSError:
        return [], ""
    candidates.sort(key=lambda c: c[0])
    images: list[dict[str, str]] = []
    total = 0
    note = ""
    for _mtime, size, p, ext in candidates:
        if len(images) >= MAX_IMAGES:
            note = f"{len(candidates) - len(images)} more image(s) not shown (limit {MAX_IMAGES})"
            break
        if size > MAX_IMAGE_BYTES:
            note = f"{p.name} too large to preview (>{MAX_IMAGE_BYTES // 1_000_000}MB)"
            continue
        if total + size > MAX_IMAGES_TOTAL_BYTES:
            note = "image preview size budget reached; some images not shown"
            break
        try:
            raw = p.read_bytes()
        except OSError:
            continue
        total += len(raw)
        b64 = base64.b64encode(raw).decode("ascii")
        images.append({
            "name": str(p.relative_to(workspace)),
            "data_url": f"data:{IMAGE_MIME_BY_EXT[ext]};base64,{b64}",
        })
    return images, note


def _artifact_snapshot(workspace: Path) -> dict[str, tuple[int, int]]:
    """Capture visible workspace files so executions can report new/changed outputs."""
    root = workspace.resolve()
    snapshot: dict[str, tuple[int, int]] = {}
    try:
        for path in root.rglob("*"):
            if not path.is_file():
                continue
            rel = path.relative_to(root)
            if set(rel.parts) & SKIP_DIRS:
                continue
            if path.suffix.lower() in ARTIFACT_JUNK_EXTS or path.name in ARTIFACT_JUNK_NAMES:
                continue
            try:
                stat = path.stat()
            except OSError:
                continue
            snapshot[str(rel)] = (stat.st_size, stat.st_mtime_ns)
    except OSError:
        pass
    return snapshot


def _changed_artifacts(workspace: Path, before: dict[str, tuple[int, int]]) -> list[str]:
    after = _artifact_snapshot(workspace)
    return sorted(path for path, metadata in after.items() if before.get(path) != metadata)


@app.post("/users/{user_id}/workspaces/{chat_id}/exec", response_model=ExecResponse)
async def exec_route(user_id: str, chat_id: str, req: ExecRequest) -> ExecResponse:
    name, workspace = _workspace(user_id, chat_id)
    artifacts_before = _artifact_snapshot(workspace)
    # timeout <= 0 means "run as long as it needs" (no limit). Otherwise it's the
    # wall-clock budget in seconds, uncapped.
    _t = int(req.timeout or 0)
    timeout = _t if _t > 0 else None
    if req.kind == "python":
        command = f"/opt/talos-sandbox-venv/bin/python -c {shlex.quote(req.code or req.command)}"
    else:
        command = req.command
    started_at = time.time()
    proc = await _start_process(user_id, chat_id, ExecuteRequest(command=command))
    try:
        await asyncio.wait_for(asyncio.shield(proc.log_task), timeout=timeout)
    except asyncio.TimeoutError:
        proc.runner.kill(force=True)
        proc.status = "killed"
        proc.exit_code = 124
        proc.finished_at = time.time()
        output, _next, _truncated = _read_process_log(proc.log_path, offset=0)
        return ExecResponse(user_id=user_id, chat_id=chat_id, linux_user=name, workspace=str(workspace), stdout=output, stderr="", exit_code=124, timed_out=True, created_artifacts=_changed_artifacts(workspace, artifacts_before))
    output, _next, _truncated = _read_process_log(proc.log_path, offset=0)
    images, image_note = _collect_new_images(workspace, started_at)
    return ExecResponse(user_id=user_id, chat_id=chat_id, linux_user=name, workspace=str(workspace), stdout=output, stderr="", exit_code=proc.exit_code or 0, images=images, image_note=image_note, created_artifacts=_changed_artifacts(workspace, artifacts_before))


def _kernel_sock_path(workspace: Path) -> Path:
    return workspace / ".talos_kernel.sock"


def _stop_kernel(user_id: str, chat_id: str) -> None:
    rec = _kernels.pop(_session_key(user_id, chat_id), None)
    if not rec:
        return
    proc = rec.get("proc")
    if proc is not None:
        try:
            os.killpg(proc.pid, signal.SIGTERM)
        except (ProcessLookupError, PermissionError, OSError):
            try:
                proc.terminate()
            except Exception:
                pass


async def _ensure_kernel(user_id: str, chat_id: str) -> dict[str, Any]:
    key = _session_key(user_id, chat_id)
    name, workspace = _workspace(user_id, chat_id)
    rec = _kernels.get(key)
    if rec and rec.get("proc") is not None and rec["proc"].poll() is None:
        return rec
    sock = _kernel_sock_path(workspace)
    try:
        if sock.exists():
            sock.unlink()
    except OSError:
        pass
    env = os.environ.copy()
    env["HOME"] = str(workspace / ".talos-home")
    env["PATH"] = f"/opt/talos-sandbox-venv/bin:{env.get('PATH', '')}"
    proc = subprocess.Popen(
        ["gosu", name, "/opt/talos-sandbox-venv/bin/python", "-c", _KERNEL_SERVER_SRC, str(sock), str(workspace)],
        cwd=str(workspace),
        env=env,
        start_new_session=True,
    )
    # Wait briefly for the socket to come up.
    for _ in range(50):
        if sock.exists():
            break
        await asyncio.sleep(0.1)
    rec = {"proc": proc, "sock": str(sock)}
    _kernels[key] = rec
    return rec


def _kernel_exec_sync(sock_path: str, code: str, timeout: int) -> dict[str, Any]:
    import socket as _socket
    import struct as _struct

    c = _socket.socket(_socket.AF_UNIX, _socket.SOCK_STREAM)
    c.settimeout(timeout if timeout and timeout > 0 else None)
    try:
        c.connect(sock_path)
        payload = json.dumps({"code": code}).encode("utf-8")
        c.sendall(_struct.pack(">I", len(payload)) + payload)

        def _recvall(n: int) -> bytes | None:
            buf = b""
            while len(buf) < n:
                d = c.recv(n - len(buf))
                if not d:
                    return None
                buf += d
            return buf

        hdr = _recvall(4)
        if not hdr:
            return {"error": "kernel: no response"}
        (ln,) = _struct.unpack(">I", hdr)
        body = _recvall(ln)
        if body is None:
            return {"error": "kernel: truncated response"}
        return json.loads(body.decode("utf-8"))
    finally:
        try:
            c.close()
        except OSError:
            pass


@app.post("/users/{user_id}/workspaces/{chat_id}/kernel/execute", response_model=ExecResponse)
async def kernel_execute_route(user_id: str, chat_id: str, req: CellRequest) -> ExecResponse:
    name, workspace = _workspace(user_id, chat_id)
    artifacts_before = _artifact_snapshot(workspace)
    rec = await _ensure_kernel(user_id, chat_id)
    if not Path(rec["sock"]).exists():
        return ExecResponse(user_id=user_id, chat_id=chat_id, linux_user=name, workspace=str(workspace), stdout="", stderr="kernel: failed to start", exit_code=1)
    started_at = time.time()
    timeout = int(req.timeout or 0)
    try:
        result = await asyncio.to_thread(_kernel_exec_sync, rec["sock"], req.code, timeout)
    except Exception as exc:
        # Socket dead/stuck — drop the kernel so the next call respawns it.
        _stop_kernel(user_id, chat_id)
        return ExecResponse(user_id=user_id, chat_id=chat_id, linux_user=name, workspace=str(workspace), stdout="", stderr=f"kernel: {exc}", exit_code=1)
    stdout = str(result.get("stdout") or "")
    stderr = str(result.get("stderr") or "")
    error = str(result.get("error") or "")
    if error:
        stderr = (stderr + "\n" + error).strip() if stderr else error
    images, image_note = _collect_new_images(workspace, started_at)
    return ExecResponse(
        user_id=user_id, chat_id=chat_id, linux_user=name, workspace=str(workspace),
        stdout=_truncate(stdout, MAX_OUTPUT_CHARS), stderr=_truncate(stderr, MAX_OUTPUT_CHARS),
        exit_code=1 if error else 0, images=images, image_note=image_note,
        created_artifacts=_changed_artifacts(workspace, artifacts_before),
    )


@app.post("/users/{user_id}/workspaces/{chat_id}/kernel/reset")
def kernel_reset_route(user_id: str, chat_id: str) -> dict[str, Any]:
    """Restart the persistent kernel (clears all in-memory state)."""
    _stop_kernel(user_id, chat_id)
    return {"ok": True}


def _extract_document_text(path: Path, ext: str) -> str:
    """Extract readable text from a binary document (PDF/Word/Excel/PowerPoint)."""
    if ext == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for i, page in enumerate(reader.pages, 1):
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            if t.strip():
                pages.append(f"--- page {i} ---\n{t}")
        return "\n\n".join(pages)
    if ext == ".docx":
        import docx
        d = docx.Document(str(path))
        lines = [p.text for p in d.paragraphs]
        for tbl in d.tables:
            for row in tbl.rows:
                lines.append("\t".join(c.text for c in row.cells))
        return "\n".join(lines)
    if ext in {".xlsx", ".xlsm"}:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        out = []
        try:
            for ws in wb.worksheets:
                out.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    out.append("\t".join("" if c is None else str(c) for c in row))
        finally:
            wb.close()
        return "\n".join(out)
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"--- slide {i} ---")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        out.append("".join(run.text for run in para.runs))
        return "\n".join(out)
    return ""


@app.post("/users/{user_id}/workspaces/{chat_id}/files/read")
def read_file_route(user_id: str, chat_id: str, req: FileReadRequest) -> dict[str, Any]:
    _name, workspace = _workspace(user_id, chat_id)
    path = _safe_path(workspace, req.path)
    # Binary documents (PDF/Office) → extract text instead of returning raw bytes.
    _ext = path.suffix.lower()
    if _ext in DOC_EXTRACT_EXTS:
        if not path.is_file():
            return {"error": f"read_file: {req.path}: not found", "exit_code": 1, "path": str(path)}
        try:
            text = _extract_document_text(path, _ext)
        except Exception as exc:
            return {"error": f"read_file: {req.path}: could not extract {_ext} ({exc})", "exit_code": 1, "path": str(path)}
        text = _truncate(text.strip() or f"[{_ext} file: no extractable text]", MAX_READ_CHARS)
        return {"output": text, "exit_code": 0, "path": str(path), "extracted": True}
    try:
        if req.offset > 0 or req.limit > 0:
            start = max(int(req.offset), 1)
            limit = max(int(req.limit), 0)
            out: list[str] = []
            budget = MAX_READ_CHARS
            count = 0
            with path.open("r", encoding="utf-8", errors="replace") as fh:
                for i, line in enumerate(fh, 1):
                    if i < start:
                        continue
                    if limit and count >= limit:
                        break
                    out.append(line)
                    count += 1
                    budget -= len(line)
                    if budget <= 0:
                        out.append(f"\n... [truncated at {MAX_READ_CHARS} chars]")
                        break
            text = "".join(out)
        else:
            text = path.read_text(encoding="utf-8", errors="replace")
            text = _truncate(text, MAX_READ_CHARS)
        return {"output": text, "exit_code": 0, "path": str(path)}
    except FileNotFoundError:
        return {"error": f"read_file: {req.path}: not found", "exit_code": 1, "path": str(path)}
    except IsADirectoryError:
        return {"error": f"read_file: {req.path}: is a directory (use ls)", "exit_code": 1, "path": str(path)}
    except OSError as exc:
        return {"error": f"read_file: {req.path}: {exc}", "exit_code": 1, "path": str(path)}


@app.post("/users/{user_id}/workspaces/{chat_id}/files/write")
def write_file_route(user_id: str, chat_id: str, req: FileWriteRequest) -> dict[str, Any]:
    name, workspace = _workspace(user_id, chat_id)
    path = _safe_path(workspace, req.path)
    old = ""
    try:
        if path.exists() and path.is_file():
            old = path.read_text(encoding="utf-8")
    except (UnicodeDecodeError, OSError):
        old = ""
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(req.content, encoding="utf-8")
        # chown the file AND any parent dirs we just created up to the workspace,
        # so the agent (running as the user) can write alongside them afterwards.
        _chown_user_chain(name, workspace, path)
    except OSError as exc:
        return {"error": f"write_file: {req.path}: {exc}", "exit_code": 1, "path": str(path)}
    result: dict[str, Any] = {"output": f"Wrote {len(req.content)} bytes to {path}", "exit_code": 0, "path": str(path)}
    diff = _diff(old, req.content, req.path)
    if diff:
        result["diff"] = diff
    return result


def _leading_ws(s: str) -> str:
    return s[: len(s) - len(s.lstrip())]


def _reindent(replacement: str, tgt_indent: str, match_indent: str) -> str:
    """Shift the replacement's indentation by the difference between the matched
    block's indent and the target's indent, so a flexibly-matched edit keeps
    correct indentation."""
    if tgt_indent == match_indent:
        return replacement
    out = []
    for line in replacement.split("\n"):
        if not line.strip():
            out.append(line)
        elif line.startswith(tgt_indent):
            out.append(match_indent + line[len(tgt_indent):])
        else:
            out.append(line)
    return "\n".join(out)


def _flexible_line_match(text: str, target: str) -> tuple[int, int] | None:
    """Find a UNIQUE whitespace-tolerant match of `target` against whole lines of
    `text` (comparing each line stripped of leading/trailing whitespace). Returns
    (start, end) char offsets, or None if not found or ambiguous. Skips CRLF text
    to avoid offset issues. This lets edit_file succeed when the model's
    indentation is slightly off, instead of failing and forcing a full rewrite."""
    if "\r" in text:
        return None
    tgt_lines = target.split("\n")
    if tgt_lines and tgt_lines[-1] == "":
        tgt_lines = tgt_lines[:-1]
    if not tgt_lines:
        return None
    tgt_stripped = [l.strip() for l in tgt_lines]
    lines = text.split("\n")
    offsets: list[int] = []
    pos = 0
    for ln in lines:
        offsets.append(pos)
        pos += len(ln) + 1
    n = len(tgt_lines)
    found: list[tuple[int, int]] = []
    for i in range(0, len(lines) - n + 1):
        if [w.strip() for w in lines[i:i + n]] == tgt_stripped:
            found.append((offsets[i], offsets[i + n - 1] + len(lines[i + n - 1])))
            if len(found) > 1:
                return None
    return found[0] if len(found) == 1 else None


def _apply_edit_chunk(text: str, chunk: dict[str, Any], idx: int) -> tuple[str | None, str | None]:
    """Apply one find/replace chunk. Returns (new_text, None) or (None, error).
    Supports optional 1-indexed inclusive line-range scoping and allow_multiple."""
    target = str(chunk.get("target", chunk.get("old_string", "")) or "")
    replacement = str(chunk.get("replacement", chunk.get("new_string", "")) or "")
    allow_multiple = bool(chunk.get("allow_multiple", chunk.get("replace_all", False)))
    if not target:
        return None, f"edit #{idx + 1}: 'target' is required"
    start_line = chunk.get("start_line")
    end_line = chunk.get("end_line")
    if start_line or end_line:
        lines = text.splitlines(keepends=True)
        s = max(int(start_line or 1), 1) - 1
        e = min(int(end_line) if end_line else len(lines), len(lines))
        if s >= e:
            return None, f"edit #{idx + 1}: invalid line range {start_line}-{end_line}"
        window = "".join(lines[s:e])
        cnt = window.count(target)
        if cnt == 0:
            return None, f"edit #{idx + 1}: target not found in lines {s + 1}-{e}"
        if cnt > 1 and not allow_multiple:
            return None, f"edit #{idx + 1}: target not unique ({cnt}) in lines {s + 1}-{e}; add context or set allow_multiple"
        new_window = window.replace(target, replacement) if allow_multiple else window.replace(target, replacement, 1)
        return "".join(lines[:s]) + new_window + "".join(lines[e:]), None
    cnt = text.count(target)
    if cnt == 0:
        # Exact match failed — try a whitespace-tolerant match so a slightly-off
        # indentation doesn't force the model to rewrite the whole file.
        span = _flexible_line_match(text, target)
        if span is None:
            return None, f"edit #{idx + 1}: target not found"
        start, end = span
        match_indent = _leading_ws(text[start:end].split("\n", 1)[0])
        tgt_indent = _leading_ws(target.split("\n", 1)[0])
        return text[:start] + _reindent(replacement, tgt_indent, match_indent) + text[end:], None
    if cnt > 1 and not allow_multiple:
        return None, f"edit #{idx + 1}: target not unique ({cnt} matches); add surrounding context or set allow_multiple"
    return (text.replace(target, replacement) if allow_multiple else text.replace(target, replacement, 1)), None


@app.post("/users/{user_id}/workspaces/{chat_id}/files/edit")
def edit_file_route(user_id: str, chat_id: str, req: FileEditRequest) -> dict[str, Any]:
    name, workspace = _workspace(user_id, chat_id)
    path = _safe_path(workspace, req.path)
    try:
        old = path.read_text(encoding="utf-8")
    except FileNotFoundError:
        return {"error": f"edit_file: {req.path}: not found (use write_file to create it)", "exit_code": 1, "path": str(path)}
    except (IsADirectoryError, UnicodeDecodeError, OSError) as exc:
        return {"error": f"edit_file: {req.path}: {exc}", "exit_code": 1, "path": str(path)}
    # Prefer the multi-chunk `edits` array; fall back to single old_string/new_string.
    edits = list(req.edits or [])
    if not edits:
        if not req.old_string:
            return {"error": "edit_file: provide `edits` (list of {target, replacement}) or old_string/new_string", "exit_code": 1, "path": str(path)}
        edits = [{"target": req.old_string, "replacement": req.new_string, "allow_multiple": req.replace_all}]
    new = old
    for i, chunk in enumerate(edits):
        new, err = _apply_edit_chunk(new, chunk, i)
        if err is not None:
            return {"error": f"edit_file: {req.path}: {err}. Read the file and match it exactly.", "exit_code": 1, "path": str(path)}
    if new == old:
        return {"error": f"edit_file: {req.path}: edits produced no change", "exit_code": 1, "path": str(path)}
    path.write_text(new, encoding="utf-8")
    _run(["chown", f"{name}:{name}", str(path)])
    n = len(edits)
    result: dict[str, Any] = {"output": f"Edited {path} ({n} edit{'s' if n != 1 else ''})", "exit_code": 0, "path": str(path)}
    diff = _diff(old, new, req.path)
    if diff:
        result["diff"] = diff
    return result


@app.post("/users/{user_id}/workspaces/{chat_id}/files/grep")
def grep_route(user_id: str, chat_id: str, req: SearchRequest) -> dict[str, Any]:
    _name, workspace = _workspace(user_id, chat_id)
    root = _safe_path(workspace, req.path or ".")
    if not req.pattern:
        return {"error": "grep: pattern is required", "exit_code": 1}
    flags = re.IGNORECASE if req.ignore_case else 0
    try:
        rx = re.compile(req.pattern, flags)
    except re.error as exc:
        return {"error": f"grep: bad pattern: {exc}", "exit_code": 1}
    max_results = max(1, min(int(req.max_results or 200), 500))
    files = [root] if root.is_file() else [p for p in root.rglob("*") if p.is_file() and not (set(p.relative_to(root).parts) & SKIP_DIRS)]
    hits: list[str] = []
    for fp in files:
        if req.glob and not fnmatch.fnmatch(fp.name, req.glob) and not fnmatch.fnmatch(str(fp.relative_to(root)), req.glob):
            continue
        try:
            with fp.open("r", encoding="utf-8", errors="strict") as fh:
                for i, line in enumerate(fh, 1):
                    if rx.search(line):
                        hits.append(f"{fp}:{i}:{line.rstrip()[:500]}")
                        if len(hits) >= max_results:
                            break
        except (UnicodeDecodeError, OSError):
            continue
        if len(hits) >= max_results:
            break
    if not hits:
        return {"output": f"No matches for {req.pattern!r} under {root}", "exit_code": 0}
    out = "\n".join(hits)
    if len(hits) >= max_results:
        out += f"\n... [capped at {max_results} matches]"
    return {"output": _truncate(out), "exit_code": 0}


@app.post("/users/{user_id}/workspaces/{chat_id}/files/glob")
def glob_route(user_id: str, chat_id: str, req: SearchRequest) -> dict[str, Any]:
    _name, workspace = _workspace(user_id, chat_id)
    root = _safe_path(workspace, req.path or ".")
    if not req.pattern:
        return {"error": "glob: pattern is required", "exit_code": 1}
    if not root.is_dir():
        return {"error": f"glob: {root}: not a directory", "exit_code": 1}
    matched: list[tuple[float, str]] = []
    try:
        for p in root.rglob(req.pattern):
            if set(p.relative_to(root).parts) & SKIP_DIRS:
                continue
            try:
                matched.append((p.stat().st_mtime, str(p)))
            except OSError:
                matched.append((0, str(p)))
    except OSError as exc:
        return {"error": f"glob: {exc}", "exit_code": 1}
    matched.sort(key=lambda item: item[0], reverse=True)
    paths = [path for _mtime, path in matched[:200]]
    if not paths:
        return {"output": f"No files matching {req.pattern!r} under {root}", "exit_code": 0}
    out = "\n".join(paths)
    if len(matched) > len(paths):
        out += f"\n... [capped at {len(paths)} files]"
    return {"output": _truncate(out), "exit_code": 0}


@app.post("/users/{user_id}/workspaces/{chat_id}/files/ls")
def ls_route(user_id: str, chat_id: str, req: ListRequest) -> dict[str, Any]:
    _name, workspace = _workspace(user_id, chat_id)
    root = _safe_path(workspace, req.path or ".")
    if not root.is_dir():
        return {"error": f"ls: {root}: not a directory", "exit_code": 1}
    rows: list[tuple[bool, str, int]] = []
    try:
        for entry in root.iterdir():
            if entry.name.startswith("."):
                continue
            is_dir = entry.is_dir()
            size = 0 if is_dir else entry.stat().st_size
            rows.append((is_dir, entry.name, size))
    except OSError as exc:
        return {"error": f"ls: {exc}", "exit_code": 1}
    rows.sort(key=lambda row: (not row[0], row[1].lower()))
    lines = [f"{root}:"]
    for is_dir, name, size in rows[:200]:
        lines.append(f"  {name}/" if is_dir else f"  {name}  ({size} B)")
    if len(rows) > 200:
        lines.append(f"  ... [{len(rows) - 200} more]")
    if not rows:
        lines.append("  (empty)")
    return {"output": _truncate("\n".join(lines)), "exit_code": 0}


def _resolve_show_image_path(name: str, workspace: Path, raw: str) -> Path | None:
    """Resolve the image path for show_image. Prefer a workspace-relative file.
    If the model saved to an absolute path outside the workspace (e.g. /tmp or its
    home — a common habit), import a copy into the workspace's output/ dir so the
    image both displays AND appears in the artifacts list. Returns the readable
    path inside the workspace, or None if it can't be found/imported."""
    # 1) Normal case: a workspace-relative (or in-workspace absolute) path.
    try:
        target = _safe_path(workspace, raw)
        if target.is_file():
            return target
    except HTTPException:
        pass
    # 2) Out-of-workspace absolute path: import it if it lives somewhere the user
    #    legitimately writes (their home subtree or /tmp) and is an image.
    cand = Path(raw)
    if not cand.is_absolute() or cand.suffix.lower() not in IMAGE_MIME_BY_EXT:
        return None
    try:
        resolved = cand.resolve()
        if not resolved.is_file():
            return None
    except OSError:
        return None
    allowed_roots = [(HOME_ROOT / name).resolve(), Path("/tmp").resolve()]
    if not any(resolved == r or r in resolved.parents for r in allowed_roots):
        return None
    dest_dir = workspace / OUTPUT_DIR_NAME
    try:
        dest_dir.mkdir(parents=True, exist_ok=True)
        dest = dest_dir / resolved.name
        shutil.copyfile(resolved, dest)
        subprocess.run(["chown", "-R", f"{name}:{name}", str(dest_dir)], capture_output=True)
    except OSError:
        return None
    return dest


@app.post("/users/{user_id}/workspaces/{chat_id}/files/image")
def read_image_route(user_id: str, chat_id: str, req: FileReadRequest) -> dict[str, Any]:
    """Read an image and return it as a base64 data URL for the `show_image` tool.
    Workspace-relative paths are read directly; absolute paths the model used (e.g.
    /tmp) are imported into the workspace's output/ dir first."""
    name, workspace = _workspace(user_id, chat_id)
    path = _resolve_show_image_path(name, workspace, (req.path or "").strip())
    if path is None:
        return {"error": f"show_image: {req.path}: not found (save the image inside your workspace, e.g. output/chart.png)", "exit_code": 1}
    ext = path.suffix.lower()
    if ext not in IMAGE_MIME_BY_EXT:
        return {"error": f"show_image: {req.path}: not a supported image type ({', '.join(sorted(IMAGE_MIME_BY_EXT))})", "exit_code": 1}
    try:
        size = path.stat().st_size
        if size > MAX_IMAGE_BYTES:
            return {"error": f"show_image: {req.path}: too large ({size // 1_000_000}MB > {MAX_IMAGE_BYTES // 1_000_000}MB limit)", "exit_code": 1}
        raw = path.read_bytes()
    except OSError as exc:
        return {"error": f"show_image: {req.path}: {exc}", "exit_code": 1}
    try:
        display_name = str(path.relative_to(workspace.resolve()))
    except ValueError:
        display_name = path.name
    b64 = base64.b64encode(raw).decode("ascii")
    return {"output": f"Displaying {display_name}", "name": display_name, "data_url": f"data:{IMAGE_MIME_BY_EXT[ext]};base64,{b64}", "exit_code": 0}


@app.get("/users/{user_id}/workspaces/{chat_id}/artifacts")
def list_artifacts_route(user_id: str, chat_id: str) -> dict[str, Any]:
    """List files in a chat's workspace (uploads, generated results, etc.) so the
    UI can show + download them. Recursive, skipping noise dirs; newest first."""
    _name, workspace = _workspace(user_id, chat_id)
    root = workspace.resolve()
    items: list[dict[str, Any]] = []
    try:
        for p in root.rglob("*"):
            if not p.is_file():
                continue
            rel = p.relative_to(root)
            if set(rel.parts) & SKIP_DIRS:
                continue
            ext = p.suffix.lower()
            # Skip obvious junk that isn't a real deliverable.
            if ext in ARTIFACT_JUNK_EXTS or p.name in ARTIFACT_JUNK_NAMES:
                continue
            try:
                st = p.stat()
            except OSError:
                continue
            items.append({
                "path": str(rel),
                "name": p.name,
                "size": st.st_size,
                "mtime": st.st_mtime,
                "mime": mimetypes.guess_type(p.name)[0] or "application/octet-stream",
                "is_image": ext in IMAGE_MIME_BY_EXT,
            })
    except OSError as exc:
        return {"artifacts": [], "error": str(exc)}
    # Sort intentional outputs (the output/ dir) first, then newest first.
    items.sort(key=lambda it: (0 if it["path"].split("/", 1)[0] == OUTPUT_DIR_NAME else 1, -it["mtime"]))
    return {"artifacts": items[:500]}


@app.get("/users/{user_id}/workspaces/{chat_id}/files/download")
def download_file_route(user_id: str, chat_id: str, path: str):
    """Stream a workspace file's raw bytes (for the UI's download/preview)."""
    from fastapi.responses import FileResponse

    _name, workspace = _workspace(user_id, chat_id)
    target = _safe_path(workspace, path)
    if not target.is_file():
        raise HTTPException(404, "File not found")
    mime = mimetypes.guess_type(target.name)[0] or "application/octet-stream"
    return FileResponse(str(target), media_type=mime, filename=target.name)


@app.get("/users/{user_id}/workspaces/{chat_id}/files/office-preview")
def office_preview_route(user_id: str, chat_id: str, path: str):
    """Convert an Office file to PDF using a real office layout engine."""
    from fastapi.responses import FileResponse

    name, workspace = _workspace(user_id, chat_id)
    target = _safe_path(workspace, path)
    if not target.is_file():
        raise HTTPException(404, "File not found")
    extension = target.suffix.lower()
    filters = {
        ".doc": "writer_pdf_Export", ".docx": "writer_pdf_Export",
        ".xls": "calc_pdf_Export", ".xlsx": "calc_pdf_Export", ".xlsm": "calc_pdf_Export",
        ".ppt": "impress_pdf_Export", ".pptx": "impress_pdf_Export",
    }
    if extension not in filters:
        raise HTTPException(400, "Only Office documents can be previewed")

    try:
        digest = hashlib.sha256(target.read_bytes()).hexdigest()
        cache_dir = workspace / ".talos-preview" / digest
        pdf_path = cache_dir / f"{target.stem}.pdf"
        if not pdf_path.is_file():
            cache_dir.mkdir(parents=True, exist_ok=True)
            _chown_user_chain(name, workspace, cache_dir)
            profile = cache_dir / f"profile-{uuid.uuid4().hex}"
            profile.mkdir()
            _chown_user_chain(name, workspace, profile)
            try:
                result = subprocess.run(
                    [
                        "gosu", name, "soffice", "--headless", "--nologo", "--nodefault",
                        "--nolockcheck", "--nofirststartwizard",
                        f"-env:UserInstallation={profile.as_uri()}",
                        "--convert-to", f"pdf:{filters[extension]}",
                        "--outdir", str(cache_dir), str(target),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=90,
                    env={
                        **os.environ,
                        "HOME": str(workspace / ".talos-home"),
                        "XDG_CACHE_HOME": str(profile / "cache"),
                        "XDG_CONFIG_HOME": str(profile / "config"),
                    },
                )
            finally:
                shutil.rmtree(profile, ignore_errors=True)
            if result.returncode != 0 or not pdf_path.is_file():
                detail = (result.stderr or result.stdout or "conversion failed").strip()
                raise HTTPException(422, detail[-500:])
    except subprocess.TimeoutExpired:
        raise HTTPException(504, "Office preview conversion timed out")
    except OSError as exc:
        raise HTTPException(500, f"Office preview conversion failed: {exc}")

    return FileResponse(str(pdf_path), media_type="application/pdf", filename=pdf_path.name)


@app.post("/users/{user_id}/workspaces/{chat_id}/files/delete")
def delete_path_route(user_id: str, chat_id: str, req: PathRequest) -> dict[str, Any]:
    _name, workspace = _workspace(user_id, chat_id)
    target = _safe_path(workspace, req.path)
    if target.resolve() == workspace.resolve():
        return {"error": "delete: refusing to delete the workspace root", "exit_code": 1}
    if not target.exists():
        return {"error": f"delete: {req.path}: not found", "exit_code": 1}
    try:
        if target.is_dir():
            shutil.rmtree(target, ignore_errors=True)
        else:
            target.unlink()
    except OSError as exc:
        return {"error": f"delete: {req.path}: {exc}", "exit_code": 1}
    return {"output": f"Deleted {req.path}", "exit_code": 0}


@app.post("/users/{user_id}/workspaces/{chat_id}/files/move")
def move_path_route(user_id: str, chat_id: str, req: MoveRequest) -> dict[str, Any]:
    name, workspace = _workspace(user_id, chat_id)
    src = _safe_path(workspace, req.src)
    dst = _safe_path(workspace, req.dst)
    if not src.exists():
        return {"error": f"move: {req.src}: not found", "exit_code": 1}
    try:
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(src), str(dst))
        _chown_user_chain(name, workspace, dst)
    except OSError as exc:
        return {"error": f"move: {exc}", "exit_code": 1}
    return {"output": f"Moved {req.src} -> {req.dst}", "exit_code": 0}


@app.post("/users/{user_id}/workspaces/{chat_id}/files/mkdir")
def mkdir_route(user_id: str, chat_id: str, req: PathRequest) -> dict[str, Any]:
    name, workspace = _workspace(user_id, chat_id)
    target = _safe_path(workspace, req.path)
    try:
        target.mkdir(parents=True, exist_ok=True)
        _chown_user_chain(name, workspace, target)
    except OSError as exc:
        return {"error": f"mkdir: {exc}", "exit_code": 1}
    return {"output": f"Created {req.path}", "exit_code": 0}


@app.get("/users/{user_id}/workspaces/{chat_id}/files/zip")
def zip_workspace_route(user_id: str, chat_id: str):
    """Stream the whole workspace as a .zip (for the UI's 'download all')."""
    import io
    import zipfile
    from fastapi.responses import Response

    _name, workspace = _workspace(user_id, chat_id)
    root = workspace.resolve()
    buf = io.BytesIO()
    total = 0
    max_bytes = 200_000_000
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in root.rglob("*"):
            if not p.is_file():
                continue
            rel = p.relative_to(root)
            if set(rel.parts) & SKIP_DIRS:
                continue
            try:
                total += p.stat().st_size
                if total > max_bytes:
                    break
                zf.write(str(p), str(rel))
            except OSError:
                continue
    buf.seek(0)
    return Response(
        content=buf.read(),
        media_type="application/zip",
        headers={"Content-Disposition": 'attachment; filename="workspace.zip"'},
    )
